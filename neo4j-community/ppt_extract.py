from pptx import Presentation
'''
from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save('test.pptx')
'''

path = '/Users/mac/major_graph/课件/'

prs = Presentation(path+'01引言_软件学院.pptx')

keys = ['title','contents','image']


#print(prs.slides)
chapter = []
for slide in prs.slides:
    #print(prs.slides.index(slide))
    #print(slide.shapes.title)
    page = []
    for shape in slide.shapes:
        shape_contents = []
        #paragraph_contents = []
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            paragraph_contents = []
            paragraph_contents.append(paragraph.text)
            paragraph_contents.append(paragraph.level)
            shape_contents.append(paragraph_contents)
        page.append(shape_contents)
    chapter.append(page)
print(chapter[12])
