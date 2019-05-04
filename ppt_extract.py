from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def get_height(lofl):
    return lofl[0][3]

def get_text(slide):
    txt = ''
    for shape in slide:
        for para in shape:
            txt = txt+para[0]
    return txt

def extract_ppt (path,file_name):

    prs = Presentation(path+file_name)

    chapter = dict()
    addition = list()
    for slide in prs.slides:
        page = []
        for shape in slide.shapes:
            shape_contents = []
            if shape.top < 6100000:
                if not shape.has_text_frame:
                    #continue
                    if prs.slides.index(slide) not in addition:
                        addition.append(prs.slides.index(slide))
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_contents = []
                    paragraph_contents.append(paragraph.text)
                    paragraph_contents.append(paragraph.level)
                    if shape.is_placeholder:
                        paragraph_contents.append(shape.name)
                    else:
                        paragraph_contents.append('')
                    paragraph_contents.append(shape.top)
                    shape_contents.append(paragraph_contents)
            page.append(shape_contents)
        chapter[str(prs.slides.index(slide))] = page

    print(chapter['94'])


    '''The structure of the chapter is the chapter[page_num][shape_num][paragraph_num], for example, 
    chapter[1][1][0] returns ['数据库系统概念----引言', 0, '页脚占位符 5', 6477000]'''

    # Delete duplicates and nulls, sort the shapes
    placeword = ['灯片编号占位符', '页脚占位符','日期占位符']

    for slide in chapter.values():
        for shape in slide:
            for paragraph in shape:
                # print(paragraph[2])
                # print(paragraph)
                if paragraph[0] == '':
                    shape.remove(paragraph)
                    # print(paragraph)
                else:
                    for word in placeword:
                        # print(word)
                        # print(paragraph[2])
                        if word in paragraph[2]:
                            shape.remove(paragraph)
                            # print(paragraph)

    previous = get_text(chapter['0'])


    for i in range(0,len(chapter.keys())):
        #print(chapter[str(i)])
        if i != 0:
            current = get_text(chapter[str(i)])
            if previous == current:
                chapter[str(i)] = []
            else:
                previous = current
        while [] in chapter[str(i)]:
            chapter[str(i)].remove([])
        if not chapter[str(i)]:
            chapter[str(i)].sort(key=get_height)

#    for slide in chapter.values():
#        slide.sort(key=get_height)

    #print(chapter)

    return chapter,addition

