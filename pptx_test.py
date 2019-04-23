from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
path = '/Users/mac/major_graph/课件/'
file_name = '01引言_软件学院.pptx'
prs = Presentation(path + file_name)
num = 0
for slide in prs.slides:
    if prs.slides.index(slide) == 38:
        for shape in slide.shapes:
            print(slide.shapes.index(shape))
            print(shape.shape_type)
            '''
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                num +=1
                #print(shape.image.blob)
                print(prs.slides.index(slide))
            '''